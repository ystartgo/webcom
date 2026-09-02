import os
import sys
import subprocess
import shutil

# Ensure working directory is always script directory
os.chdir(os.path.dirname(os.path.abspath(__file__)))
if os.path.dirname(os.path.abspath(__file__)) not in sys.path:
    sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from typing import Optional
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
import uvicorn

app = FastAPI(title="Webcom Multi-Protocol Backend Daemon", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

models_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "models")
os.makedirs(models_dir, exist_ok=True)

from fastapi.responses import FileResponse

@app.get("/models/{model_name}/resolve/main/{filepath:path}")
def serve_model_resolve_main(model_name: str, filepath: str):
    """相容 WebLLM 請求 HuggingFace resolve/main 格式路徑"""
    target = os.path.join(models_dir, model_name, filepath)
    if os.path.exists(target) and os.path.isfile(target):
        return FileResponse(target)
    if filepath == "tensor-cache.json":
        fallback = os.path.join(models_dir, model_name, "ndarray-cache.json")
        if os.path.exists(fallback) and os.path.isfile(fallback):
            return FileResponse(fallback)
    raise HTTPException(status_code=404, detail=f"File not found: {filepath}")


app.mount("/models", StaticFiles(directory=models_dir), name="models")


# ── 根目錄靜態服務 (index.html 透過 http://127.0.0.1:8001 開啟，解決 file:// Cache API 限制) ──
webcom_dir = os.path.dirname(os.path.abspath(__file__))

@app.get("/")
def serve_index():
    """提供 Webcom index.html 根頁面，讓瀏覽器可以透過 HTTP 協議存取，解決 WebGPU Cache API 在 file:// 下被封鎖的問題"""
    from fastapi.responses import FileResponse
    index_path = os.path.join(webcom_dir, "index.html")
    if os.path.exists(index_path):
        return FileResponse(index_path, media_type="text/html")
    raise HTTPException(status_code=404, detail="index.html not found")

# 掛載 assets 子目錄（JS/圖片等靜態資源）
assets_dir = os.path.join(webcom_dir, "assets")
if os.path.exists(assets_dir):
    app.mount("/assets", StaticFiles(directory=assets_dir), name="assets")

@app.get("/api/local_models")
def get_local_models():
    """Detect and return locally downloaded WebLLM model folders"""
    installed = []
    if os.path.exists(models_dir):
        for item in os.listdir(models_dir):
            item_path = os.path.join(models_dir, item)
            if os.path.isdir(item_path):
                # Check for mlc-chat-config.json or ndarray-cache.json
                has_config = os.path.exists(os.path.join(item_path, "mlc-chat-config.json")) or os.path.exists(os.path.join(item_path, "ndarray-cache.json"))
                size_mb = 0
                try:
                    total_bytes = sum(os.path.getsize(os.path.join(item_path, f)) for f in os.listdir(item_path) if os.path.isfile(os.path.join(item_path, f)))
                    size_mb = round(total_bytes / (1024 * 1024), 1)
                except Exception:
                    pass
                installed.append({
                    "id": item,
                    "name": item,
                    "ready": has_config,
                    "size_mb": size_mb,
                    "url": f"http://127.0.0.1:8001/models/{item}/"
                })
    return {"models": installed}

import logging
import threading
import urllib.request
import json

downloading_models = set()

def _bg_download_model_worker(model_name: str):
    try:
        logging.info(f"Auto-caching model '{model_name}' to local disk...")
        target_dir = os.path.join(models_dir, model_name)
        os.makedirs(target_dir, exist_ok=True)
        
        if os.path.exists(os.path.join(target_dir, "ndarray-cache.json")):
            has_shards = any(f.startswith("params_shard") for f in os.listdir(target_dir))
            if has_shards:
                logging.info(f"Model '{model_name}' already complete on disk.")
                downloading_models.discard(model_name)
                return

        hf_base = f"https://huggingface.co/mlc-ai/{model_name}/resolve/main"
        
        cache_json_url = f"{hf_base}/ndarray-cache.json"
        cache_json_path = os.path.join(target_dir, "ndarray-cache.json")
        req = urllib.request.Request(cache_json_url, headers={'User-Agent': 'Webcom-AutoCache/1.0'})
        with urllib.request.urlopen(req, timeout=20) as resp, open(cache_json_path, 'wb') as f:
            f.write(resp.read())

        for fn in ["mlc-chat-config.json", "tokenizer.json", "tokenizer_config.json", "vocab.json", "merges.txt"]:
            try:
                url = f"{hf_base}/{fn}"
                dest = os.path.join(target_dir, fn)
                req = urllib.request.Request(url, headers={'User-Agent': 'Webcom-AutoCache/1.0'})
                with urllib.request.urlopen(req, timeout=20) as resp, open(dest, 'wb') as f:
                    f.write(resp.read())
            except Exception:
                pass

        with open(cache_json_path, "r", encoding="utf-8") as f:
            cache_data = json.load(f)
        records = cache_data.get("records", [])
        for r in records:
            sf = r.get("dataPath")
            if sf:
                shard_url = f"{hf_base}/{sf}"
                shard_dest = os.path.join(target_dir, sf)
                req = urllib.request.Request(shard_url, headers={'User-Agent': 'Webcom-AutoCache/1.0'})
                with urllib.request.urlopen(req, timeout=60) as resp, open(shard_dest, 'wb') as f:
                    f.write(resp.read())

        wasm_url = "https://raw.githubusercontent.com/mlc-ai/binary-mlc-llm-libs/main/web-llm-models/v0_2_48/Qwen2-0.5B-Instruct-q4f16_1-ctx4k_cs1k-webgpu.wasm"
        wasm_dest = os.path.join(target_dir, "Qwen2-0.5B-Instruct-q4f16_1-ctx4k_cs1k-webgpu.wasm")
        try:
            req = urllib.request.Request(wasm_url, headers={'User-Agent': 'Webcom-AutoCache/1.0'})
            with urllib.request.urlopen(req, timeout=30) as resp, open(wasm_dest, 'wb') as f:
                f.write(resp.read())
        except Exception:
            pass

        logging.info(f"✅ Model '{model_name}' successfully cached to local disk: {target_dir}!")
    except Exception as e:
        logging.error(f"Failed to auto-cache model '{model_name}': {e}")
    finally:
        downloading_models.discard(model_name)

class AutoCacheRequest(BaseModel):
    model_name: Optional[str] = "Qwen2.5-0.5B-Instruct-q4f16_1-MLC"

@app.post("/api/auto_cache_model")
def auto_cache_model(req: AutoCacheRequest):
    model_name = req.model_name or "Qwen2.5-0.5B-Instruct-q4f16_1-MLC"
    if model_name in downloading_models:
        return {"status": "in_progress", "model": model_name}
    
    target_dir = os.path.join(models_dir, model_name)
    if os.path.exists(os.path.join(target_dir, "ndarray-cache.json")) and any(f.startswith("params_shard") for f in os.listdir(target_dir)):
        return {"status": "already_cached", "model": model_name}

    downloading_models.add(model_name)
    t = threading.Thread(target=_bg_download_model_worker, args=(model_name,), daemon=True)
    t.start()
    return {"status": "started", "model": model_name}

log_file_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "daemon.log")
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)]
)
logging.info("==================================================")
logging.info("Webcom Daemon server process initialized on port 8001")
logging.info("==================================================")

@app.get("/api/logs")
def get_daemon_logs(lines: int = 100):
    """Return latest lines from daemon.log for LLM diagnostics"""
    if not os.path.exists(log_file_path):
        return {"status": "success", "logs": "No logs recorded yet."}
    try:
        with open(log_file_path, "r", encoding="utf-8", errors="replace") as f:
            all_lines = f.readlines()
            tail = "".join(all_lines[-lines:])
            return {"status": "success", "total_lines": len(all_lines), "logs": tail}
    except Exception as e:
        return {"status": "error", "error": str(e)}

class ShellRequest(BaseModel):
    command: str

class FileReadRequest(BaseModel):
    filepath: str

class FileWriteRequest(BaseModel):
    filepath: str
    content: str

class SSHRequest(BaseModel):
    host: str
    username: str
    password: Optional[str] = ""
    command: str
    port: Optional[int] = 22

class TelnetRequest(BaseModel):
    host: str
    port: Optional[int] = 23
    username: Optional[str] = ""
    password: Optional[str] = ""
    command: str
    cmd_prompt: Optional[str] = "#"

class SerialRequest(BaseModel):
    port: str
    baudrate: Optional[int] = 115200
    command: str

@app.get("/")
@app.get("/health")
def health_check():
    return {"status": "online", "service": "Webcom Daemon", "port": 8001}

@app.post("/shutdown")
def shutdown_daemon():
    import threading, time
    def delayed_exit():
        time.sleep(0.5)
        os._exit(0)
    threading.Thread(target=delayed_exit).start()
    return {"status": "success", "message": "Daemon shutting down"}

@app.post("/tools/execute_shell")
def execute_shell(req: ShellRequest):
    cmd = req.command.strip()
    if not cmd:
        return {"status": "success", "stdout": "", "stderr": ""}
    
    try:
        is_windows = sys.platform.startswith("win")
        process = None

        if is_windows:
            # WinPE & Windows Execution: Prefer powershell with UTF-8 console encoding if available, fallback cleanly to cmd.exe
            if shutil.which("powershell.exe"):
                try:
                    process = subprocess.run(
                        ["powershell.exe", "-NoProfile", "-Command", f"[Console]::OutputEncoding = [System.Text.Encoding]::UTF8; {cmd}"],
                        capture_output=True,
                        text=True,
                        timeout=30,
                        encoding="utf-8",
                        errors="replace"
                    )
                except Exception:
                    process = None

            if process is None:
                # Native CMD execution (100% available in all WinPE builds)
                process = subprocess.run(
                    f'cmd.exe /c "{cmd}"',
                    shell=True,
                    capture_output=True,
                    text=True,
                    timeout=30,
                    encoding="utf-8",
                    errors="replace"
                )
        else:
            process = subprocess.run(
                cmd,
                shell=True,
                capture_output=True,
                text=True,
                timeout=30,
                encoding="utf-8",
                errors="replace"
            )
        
        return {
            "status": "success" if process.returncode == 0 else "error",
            "returncode": process.returncode,
            "stdout": process.stdout,
            "stderr": process.stderr
        }
    except subprocess.TimeoutExpired:
        return {"status": "error", "error": "指令執行超時 (Timeout 30s)"}
    except Exception as e:
        return {"status": "error", "error": str(e)}

class SearchRequest(BaseModel):
    query: str

@app.post("/tools/web_search")
def web_search_endpoint(req: SearchRequest):
    query = req.query.strip()
    if not query:
        return {"status": "error", "error": "搜尋關鍵字不得為空"}
    
    # 優先處理 IP / GEO 地理位置查詢
    if any(k in query.lower() for k in ["ip", "geo", "地理位置", "經緯度", "所在城市", "定位", "電信業者", "isp"]):
        try:
            req_geo = urllib.request.Request("http://ip-api.com/json/?fields=status,message,country,countryCode,region,regionName,city,zip,lat,lon,timezone,isp,org,as,query", headers={'User-Agent': 'curl/7.68.0'})
            with urllib.request.urlopen(req_geo, timeout=6) as resp:
                geo_data = json.loads(resp.read().decode('utf-8', errors='replace'))
                if geo_data.get("status") == "success":
                    info = (
                        f"【本機外網 IP 與 GEO 地理位置資訊】\n"
                        f"- 外網 IP: {geo_data.get('query')}\n"
                        f"- 國家/地區: {geo_data.get('country')} ({geo_data.get('countryCode')})\n"
                        f"- 所在城市/區域: {geo_data.get('city')}, {geo_data.get('regionName')}\n"
                        f"- 經緯度座標: Lat {geo_data.get('lat')}, Lon {geo_data.get('lon')}\n"
                        f"- 系統時區: {geo_data.get('timezone')}\n"
                        f"- 網際網路供應商 (ISP): {geo_data.get('isp')} ({geo_data.get('org')})"
                    )
                    return {"status": "success", "query": query, "result": info}
        except Exception as e:
            logging.warning(f"IP Geo lookup failed: {e}")

    # 優先處理停班停課 (颱風假) 查詢
    if any(k in query.lower() for k in ["停班停課", "颱風假", "上班上課", "人事行政總處", "天災假", "停班", "停課"]):
        try:
            # 優先搜尋行政院人事行政總處與氣象最新通報
            dgpa_search_url = f"https://html.duckduckgo.com/html/?q={urllib.parse.quote('行政院人事行政總處 各縣市 停班停課 最新公告')}"
            headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'}
            req_dgpa = urllib.request.Request(dgpa_search_url, headers=headers)
            with urllib.request.urlopen(req_dgpa, timeout=8) as resp:
                html = resp.read().decode('utf-8', errors='replace')
                import re
                snippets = re.findall(r'<a class="result__snippet[^>]*>(.*?)</a>', html, re.DOTALL)
                clean_snippets = [re.sub(r'<[^>]+>', '', s).strip() for s in snippets[:3] if re.sub(r'<[^>]+>', '', s).strip()]
                if clean_snippets:
                    return {"status": "success", "query": query, "result": f"【行政院人事行政總處與天然災害停班停課最新通報】\n" + "\n---\n".join(clean_snippets)}
        except Exception as e:
            logging.warning(f"DGPA search error: {e}")

    # 優先處理天氣查詢 (透過 wttr.in)
    if any(k in query.lower() for k in ["天氣", "weather", "溫度", "氣溫", "降雨", "氣象"]):
        try:
            clean_q = query.replace("天氣", "").replace("即時", "").replace("查詢", "").strip() or "Taipei"
            req_weather = urllib.request.Request(f"https://wttr.in/{urllib.parse.quote(clean_q)}?format=%l:+%c+%t,+濕度:%h,+風速:%w,+%p&lang=zh-tw", headers={'User-Agent': 'curl/7.68.0'})
            with urllib.request.urlopen(req_weather, timeout=6) as resp:
                weather_data = resp.read().decode('utf-8', errors='replace').strip()
                if weather_data and "404" not in weather_data and "Unknown" not in weather_data:
                    return {"status": "success", "query": query, "result": f"【即時氣象資訊 (wttr.in)】\n{weather_data}"}
        except Exception:
            pass

    # DuckDuckGo HTML 網頁搜尋 (後端執行無 CORS 限制)
    try:
        search_url = f"https://html.duckduckgo.com/html/?q={urllib.parse.quote(query)}"
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'}
        req_obj = urllib.request.Request(search_url, headers=headers)
        with urllib.request.urlopen(req_obj, timeout=8) as resp:
            html = resp.read().decode('utf-8', errors='replace')
            import re
            snippets = re.findall(r'<a class="result__snippet[^>]*>(.*?)</a>', html, re.DOTALL)
            clean_snippets = []
            for snip in snippets[:4]:
                clean_text = re.sub(r'<[^>]+>', '', snip).strip()
                if clean_text:
                    clean_snippets.append(clean_text)
            if clean_snippets:
                return {"status": "success", "query": query, "result": "\n---\n".join(clean_snippets)}
    except Exception as e:
        logging.warning(f"DuckDuckGo search error: {e}")

    # 維基百科 API 備援
    try:
        wiki_url = f"https://zh.wikipedia.org/w/api.php?action=query&list=search&srsearch={urllib.parse.quote(query)}&format=json"
        req_wiki = urllib.request.Request(wiki_url, headers={'User-Agent': 'Webcom/1.0'})
        with urllib.request.urlopen(req_wiki, timeout=6) as resp:
            wiki_json = json.loads(resp.read().decode('utf-8', errors='replace'))
            items = wiki_json.get("query", {}).get("search", [])
            if items:
                import re
                results = [f"【{it.get('title')}】 {re.sub(r'<[^>]+>', '', it.get('snippet', ''))}" for it in items[:3]]
                return {"status": "success", "query": query, "result": "\n---\n".join(results)}
    except Exception as e:
        logging.warning(f"Wikipedia search error: {e}")

    return {"status": "error", "error": f"搜尋失敗或網路無法連線"}

# ─────────────────────────────────────────────────────────────
# 📄 Document Parser Endpoint (PDF / Word / Excel / PPT → Text)
# ─────────────────────────────────────────────────────────────
import base64
import tempfile

class DocumentParseRequest(BaseModel):
    filename: str       # 原始檔名 (判斷副檔名用)
    data_base64: str    # 前端 FileReader.readAsDataURL → base64 部分

@app.post("/tools/parse_document")
def parse_document(req: DocumentParseRequest):
    """
    接收 Base64 編碼的文件，根據副檔名使用對應解析器轉換為純文字。
    支援: .pdf, .docx, .doc, .xlsx, .xls, .pptx, .ppt, .txt, .md, .csv, .json
    """
    ext = os.path.splitext(req.filename.lower())[1]
    try:
        # 解碼 base64 (支援 data:...;base64,<data> 或純 base64)
        raw_b64 = req.data_base64
        if "base64," in raw_b64:
            raw_b64 = raw_b64.split("base64,", 1)[1]
        file_bytes = base64.b64decode(raw_b64)
    except Exception as e:
        return {"status": "error", "error": f"Base64 解碼失敗: {e}"}

    text = ""
    try:
        # ── PDF ──────────────────────────────────────────────
        if ext == ".pdf":
            try:
                import io, pypdf
                reader = pypdf.PdfReader(io.BytesIO(file_bytes))
                text = "\n".join(p.extract_text() or "" for p in reader.pages)
            except Exception as e:
                try:
                    import io, pdfminer.high_level as pm
                    text = pm.extract_text(io.BytesIO(file_bytes))
                except Exception as e2:
                    return {"status": "error", "error": f"PDF 解析失敗: {e}"}


        # ── Word (.docx) ──────────────────────────────────────
        elif ext in (".docx",):
            try:
                import docx
                import io
                doc = docx.Document(io.BytesIO(file_bytes))
                parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
                for tbl in doc.tables:
                    for row in tbl.rows:
                        parts.append("\t".join(cell.text for cell in row.cells))
                text = "\n".join(parts)
            except ImportError:
                return {"status": "error", "error": "解析 .docx 需要 python-docx 套件。\n請執行: pip install python-docx"}

        # ── Word (.doc 舊格式) ─────────────────────────────────
        elif ext == ".doc":
            try:
                with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                result = subprocess.run(
                    ["antiword", tmp_path], capture_output=True, text=True, timeout=15
                )
                os.unlink(tmp_path)
                text = result.stdout or result.stderr
                if not text.strip():
                    return {"status": "error", "error": "舊版 .doc 格式請安裝 antiword 或改用 .docx 格式"}
            except FileNotFoundError:
                return {"status": "error", "error": "解析舊版 .doc 需要 antiword 工具，或請將文件另存為 .docx"}

        # ── Excel ─────────────────────────────────────────────
        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                import io
                wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
                parts = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    parts.append(f"=== 工作表: {sheet_name} ===")
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        if any(c.strip() for c in cells):
                            parts.append("\t".join(cells))
                text = "\n".join(parts)
            except ImportError:
                return {"status": "error", "error": "解析 Excel 需要 openpyxl 套件。\n請執行: pip install openpyxl"}

        # ── PowerPoint ────────────────────────────────────────
        elif ext in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                import io
                prs = Presentation(io.BytesIO(file_bytes))
                parts = []
                for i, slide in enumerate(prs.slides, 1):
                    parts.append(f"=== 投影片 {i} ===")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text)
                text = "\n".join(parts)
            except ImportError:
                return {"status": "error", "error": "解析 PowerPoint 需要 python-pptx 套件。\n請執行: pip install python-pptx"}

        # ── 純文字、網頁與各類原始碼程式檔案 ───────
        elif ext in (
            ".txt", ".md", ".markdown", ".csv", ".tsv", ".log", ".json", ".jsonl", ".svg", ".xml",
            ".html", ".htm", ".xhtml", ".js", ".jsx", ".ts", ".tsx", ".css", ".scss", ".less",
            ".py", ".bat", ".cmd", ".ps1", ".sh", ".bash", ".yaml", ".yml", ".ini", ".env",
            ".toml", ".conf", ".config", ".sql", ".c", ".cpp", ".h", ".hpp", ".cs", ".go",
            ".rs", ".java", ".kt", ".php", ".vue", ".svelte", ".r", ".rb", ".dockerfile", ".makefile"
        ):
            text = file_bytes.decode("utf-8", errors="replace")

        else:
            # 嘗試作為通用 UTF-8 純文字讀取，若無亂碼則直接支援
            try:
                decoded = file_bytes.decode("utf-8")
                text = decoded
            except Exception:
                return {"status": "error", "error": f"不支援的檔案格式: {ext}。支援格式: PDF, DOCX, XLSX, PPTX, SVG, HTML, JS, CSS, TXT, MD, CSV, JSON 等純文字與代碼檔案"}

        text = (text or "").strip()
        if not text:
            return {"status": "error", "error": "文件解析成功，但未能提取到任何文字內容（可能為掃描圖像 PDF 或空檔案）"}

        char_count = len(text)
        return {
            "status": "success",
            "filename": req.filename,
            "ext": ext,
            "char_count": char_count,
            "text": text  # 支援完整大文字/原始碼檔案 (包含 index.html 等完整內容)
        }

    except Exception as e:
        logging.exception(f"Document parse error ({req.filename}): {e}")
        return {"status": "error", "error": f"文件解析異常: {str(e)}"}


@app.get("/api/rag/list")
def list_rag_assets():
    """自動掃描 assets/RAG/ 目錄下的所有知識庫文件與範本 (支援 SVG、純文字、Markdown 等)"""
    rag_root = os.path.join(webcom_dir, "assets", "RAG")
    if not os.path.exists(rag_root):
        os.makedirs(rag_root, exist_ok=True)
    
    docs = []
    for root, dirs, files in os.walk(rag_root):
        for f in files:
            full_path = os.path.join(root, f)
            rel_path = os.path.relpath(full_path, rag_root)
            category = os.path.basename(root) if root != rag_root else "General"
            ext = os.path.splitext(f)[1].lower()
            try:
                if ext in (".svg", ".txt", ".md", ".json", ".csv", ".xml", ".py", ".yaml", ".yml"):
                    with open(full_path, "r", encoding="utf-8", errors="replace") as fh:
                        content = fh.read()
                    docs.append({
                        "id": f"rag_{rel_path.replace(os.sep, '_')}",
                        "title": f"{category}: {f}",
                        "filename": f,
                        "rel_path": rel_path.replace("\\", "/"),
                        "category": category,
                        "content": content,
                        "ext": ext
                    })
            except Exception as ex:
                logging.warning(f"Failed to read RAG file {full_path}: {ex}")
    return {"status": "success", "count": len(docs), "docs": docs}




@app.post("/tools/read_file")
def read_file(req: FileReadRequest):
    try:
        if not os.path.exists(req.filepath):
            return {"status": "error", "error": f"檔案不存在: {req.filepath}"}
        with open(req.filepath, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        return {"status": "success", "content": content}
    except Exception as e:
        return {"status": "error", "error": str(e)}

@app.post("/tools/write_file")
def write_file(req: FileWriteRequest):
    try:
        os.makedirs(os.path.dirname(os.path.abspath(req.filepath)), exist_ok=True)
        with open(req.filepath, "w", encoding="utf-8") as f:
            f.write(req.content)
        return {"status": "success", "message": f"成功寫入檔案: {req.filepath}"}
    except Exception as e:
        return {"status": "error", "error": str(e)}

@app.post("/tools/execute_ssh")
def execute_ssh(req: SSHRequest):
    try:
        import paramiko
        client = paramiko.SSHClient()
        client.set_missing_host_key_policy(paramiko.AutoAddPolicy())
        client.connect(hostname=req.host, port=req.port, username=req.username, password=req.password, timeout=10)
        stdin, stdout, stderr = client.exec_command(req.command, timeout=20)
        out = stdout.read().decode('utf-8', errors='replace')
        err = stderr.read().decode('utf-8', errors='replace')
        client.close()
        return {"status": "success", "stdout": out, "stderr": err}
    except ImportError:
        return {"status": "error", "error": "後端尚未安裝 paramiko 套件 (pip install paramiko)"}
    except Exception as e:
        return {"status": "error", "error": f"SSH 連線失敗: {str(e)}"}

@app.post("/tools/execute_telnet")
def execute_telnet(req: TelnetRequest):
    try:
        import telnetlib
        tn = telnetlib.Telnet(req.host, req.port, timeout=10)
        if req.username:
            tn.read_until(b"login: ", timeout=5)
            tn.write(req.username.encode('ascii') + b"\n")
        if req.password:
            tn.read_until(b"Password: ", timeout=5)
            tn.write(req.password.encode('ascii') + b"\n")
        
        prompt = req.cmd_prompt.encode('ascii')
        tn.read_until(prompt, timeout=5)
        tn.write(req.command.encode('utf-8') + b"\n")
        out = tn.read_until(prompt, timeout=10).decode('utf-8', errors='replace')
        tn.close()
        return {"status": "success", "stdout": out}
    except Exception as e:
        return {"status": "error", "error": f"Telnet 連線失敗: {str(e)}"}

@app.post("/tools/execute_serial")
def execute_serial(req: SerialRequest):
    try:
        import serial
        import time
        ser = serial.Serial(req.port, req.baudrate, timeout=3)
        ser.write((req.command + "\r\n").encode('utf-8'))
        time.sleep(0.5)
        out = ser.read_all().decode('utf-8', errors='replace')
        ser.close()
        return {"status": "success", "stdout": out}
    except ImportError:
        return {"status": "error", "error": "後端尚未安裝 pyserial 套件 (pip install pyserial)"}
    except Exception as e:
        return {"status": "error", "error": f"Serial 操作失敗: {str(e)}"}

class MCPCallRequest(BaseModel):
    name: str
    arguments: Optional[dict] = {}

@app.get("/mcp/tools/list")
def list_mcp_tools():
    return {
        "tools": [
            {
                "name": "mcp_get_system_info",
                "description": "獲取本機系統資源資訊 (CPU, Memory, Disk, Platform)",
                "inputSchema": {
                    "type": "object",
                    "properties": {}
                }
            },
            {
                "name": "mcp_list_processes",
                "description": "列出本機正在運行的重要系統行程",
                "inputSchema": {
                    "type": "object",
                    "properties": {
                        "filter_name": {"type": "string", "description": "行程名稱過濾條件 (選填)"}
                    }
                }
            },
            {
                "name": "mcp_list_directory",
                "description": "列出指定目錄下的所有檔案與資料夾",
                "inputSchema": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "目錄絕對路徑或相對路徑"}
                    },
                    "required": ["path"]
                }
            }
        ]
    }

@app.post("/mcp/tools/call")
def call_mcp_tool(req: MCPCallRequest):
    tool_name = req.name
    args = req.arguments or {}
    
    if tool_name == "mcp_get_system_info":
        import platform
        import psutil
        try:
            mem = psutil.virtual_memory()
            disk = psutil.disk_usage('/')
            return {
                "content": [{
                    "type": "text",
                    "text": f"系統: {platform.system()} {platform.release()} ({platform.machine()})\n"
                            f"CPU 核心: {psutil.cpu_count(logical=True)} (使用率: {psutil.cpu_percent()}%)\n"
                            f"記憶體: 已用 {mem.used // (1024**2)}MB / 總共 {mem.total // (1024**2)}MB ({mem.percent}%)\n"
                            f"硬碟: 已用 {disk.used // (1024**3)}GB / 總共 {disk.total // (1024**3)}GB ({disk.percent}%)"
                }]
            }
        except ImportError:
            import platform
            return {
                "content": [{
                    "type": "text",
                    "text": f"系統: {platform.system()} {platform.release()} ({platform.machine()})\n(未安裝 psutil 套件以提供詳細資源資訊)"
                }]
            }
    
    if tool_name == "mcp_list_directory":
        dir_path = args.get("path", ".")
        try:
            if not os.path.exists(dir_path):
                return {"isError": True, "content": [{"type": "text", "text": f"目錄不存在: {dir_path}"}]}
            entries = os.listdir(dir_path)
            details = []
            for item in entries[:100]:
                full = os.path.join(dir_path, item)
                is_dir = os.path.isdir(full)
                sz = os.path.getsize(full) if not is_dir else 0
                details.append(f"[{ 'DIR' if is_dir else 'FILE' }] {item} ({sz} bytes)")
            return {
                "content": [{
                    "type": "text",
                    "text": f"目錄 {dir_path} 清單 ({len(entries)} 項目):\n" + "\n".join(details)
                }]
            }
        except Exception as e:
            return {"isError": True, "content": [{"type": "text", "text": str(e)}]}

    if tool_name == "mcp_list_processes":
        import psutil
        try:
            filter_str = (args.get("filter_name") or "").lower()
            procs = []
            for p in psutil.process_iter(['pid', 'name', 'cpu_percent', 'memory_percent']):
                try:
                    pinfo = p.info
                    name = pinfo['name'] or ''
                    if filter_str and filter_str not in name.lower():
                        continue
                    procs.append(f"PID {pinfo['pid']:6d} | {name:25s} | CPU {pinfo.get('cpu_percent', 0.0):4.1f}% | MEM {pinfo.get('memory_percent', 0.0):4.1f}%")
                except (psutil.NoSuchProcess, psutil.AccessDenied):
                    pass
            return {
                "content": [{
                    "type": "text",
                    "text": "行程清單 (前 30 項):\n" + "\n".join(procs[:30])
                }]
            }
        except ImportError:
            return {"isError": True, "content": [{"type": "text", "text": "未安裝 psutil 套件"}]}

    return {"isError": True, "content": [{"type": "text", "text": f"未知的 MCP 工具: {tool_name}"}]}

if __name__ == "__main__":
    print("Webcom Daemon 正在啟動於 http://127.0.0.1:8001 ...")
    uvicorn.run(app, host="127.0.0.1", port=8001)

